import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# Function to add a background image from a URL to a slide
def add_background_image(slide, image_url, prs):
    try:
        response = requests.get(image_url)
        response.raise_for_status()
        img_data = BytesIO(response.content)
        # Add the image covering the entire slide
        background_pic = slide.shapes.add_picture(img_data, 0, 0, width=prs.slide_width, height=prs.slide_height)
        # Reorder the picture element so that it appears at the back
        sp = slide.shapes._sp
        bg_elem = background_pic._element
        sp.remove(bg_elem)
        sp.insert(0, bg_elem)
    except Exception as e:
        print(f"Warning: Could not load image from '{image_url}'. Error: {e}")

# Create a presentation object
prs = Presentation()

# URL of your lawsuit-related background image. Replace this URL with your desired image link.
background_image_url = "https://www.example.com/lawsuit_background.jpg"

# -------- Slide 1: Title Slide --------
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Legal Brief Argument Counter Argument Linking"
subtitle.text = "Law Integration with AI\nEnhancing Legal Research Through Automation"

# -------- Slide 2: Problem Statement --------
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Problem Statement"
content.text = (
    "Legal briefs contain complex arguments and counter-arguments, making manual review time-consuming and error-prone.\n\n"
    "Our solution automatically extracts and links argument sections from moving and response briefs to clarify contested points and streamline research."
)

# -------- Slide 3: Business Use Case --------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Business Use Case"
content.text = (
    "• Increased Efficiency: Automates the identification and matching of arguments.\n"
    "• Enhanced Accuracy: Minimizes manual errors in linking counter-arguments.\n"
    "• Improved Strategy: Helps attorneys quickly identify opposing arguments.\n"
    "• Better Decision-Making: Supports judges with a clear, visual overview of debates."
)

# -------- Slide 4: Use Cases Overview --------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Use Cases & Solutions"
content.text = (
    "1. Rapid Legal Research and Argument Analysis\n"
    "   - Auto-extraction of arguments and counter-arguments.\n\n"
    "2. Efficient Case Preparation\n"
    "   - Direct mapping of moving brief arguments to response brief counters.\n\n"
    "3. Judicial Decision-Making Support\n"
    "   - Visual flow of argument interconnections.\n\n"
    "4. Academic and Legal Research\n"
    "   - Structured data for trend analysis in legal debates."
)

# -------- Slide 5: Detailed Use Case 1 --------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content_box = slide.placeholders[1]
title.text = "Use Case 1: Rapid Legal Research"
content_box.text = (
    "Law firms under tight deadlines can quickly analyze briefs.\n\n"
    "Solution: The system extracts key arguments and matches them with counter-arguments using NLP techniques. "
    "Visual highlights and filters allow fast identification of contested points."
)

# -------- Slide 6: Detailed Use Case 2 --------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content_box = slide.placeholders[1]
title.text = "Use Case 2: Efficient Case Preparation"
content_box.text = (
    "Attorneys prepare cases by understanding both their own arguments and the opponent's counters.\n\n"
    "Solution: The mapping tool highlights weaknesses and opportunities, aiding in strategic rebuttals and improved case preparation."
)

# -------- Slide 7: Detailed Use Case 3 --------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content_box = slide.placeholders[1]
title.text = "Use Case 3: Judicial Decision-Making"
content_box.text = (
    "Judges require clarity in understanding argument flows to make informed decisions.\n\n"
    "Solution: The tool visually connects moving brief arguments with response counters, reducing cognitive load and ensuring key points are considered."
)

# -------- Slide 8: Detailed Use Case 4 --------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content_box = slide.placeholders[1]
title.text = "Use Case 4: Academic & Legal Research"
content_box.text = (
    "Researchers analyze trends in legal argumentation over time.\n\n"
    "Solution: Structured data from the tool can be aggregated for in-depth analysis of legal discourse and argument effectiveness."
)

# -------- Slide 9: Interactive Demo & Resources --------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content_box = slide.placeholders[1]
title.text = "Interactive Demo & Resources"
content_box.text = (
    "Explore a live demo of the system and learn more about the project.\n\n"
    "Click the link below for further details and code examples."
)

# Add a clickable hyperlink (simulate an interactive element)
left = Inches(1)
top = Inches(4)
width = Inches(5)
height = Inches(0.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Learn More: Bloomberg Hackathon - Legal Brief Linking"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# Set the hyperlink URL for the paragraph text (replace with actual URL if available)
r = p.runs[0]
r.hyperlink.address = "https://www.bloomberg.com"

# -------- Slide 10: Conclusion --------
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
add_background_image(slide, background_image_url, prs)

title = slide.shapes.title
content_box = slide.placeholders[1]
title.text = "Conclusion"
content_box.text = (
    "The Legal Brief Argument Counter Argument Linking tool automates complex legal analysis.\n\n"
    "By integrating AI with legal research, it not only streamlines the process but also enhances the precision of legal argument mapping.\n\n"
    "This innovation paves the way for smarter, faster, and more effective legal proceedings."
)

# Save the presentation
prs.save("Legal_Brief_Argument_Counter_Argument_Linking.pptx")
print("Presentation created successfully!")
